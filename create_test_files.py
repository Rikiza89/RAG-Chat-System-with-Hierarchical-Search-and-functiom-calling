#!/usr/bin/env python3
"""
Create sample office files for testing functions
"""

import os

def create_test_excel():
    """Create sample Excel file"""
    try:
        import pandas as pd
    except ImportError:
        print("❌ pandas not installed")
        return
    
    data = {
        'Product': ['Laptop', 'Mouse', 'Keyboard', 'Monitor', 'Headset'],
        'Sales': [50, 150, 100, 75, 120],
        'Revenue': [50000, 3000, 8000, 22500, 9600],
        'Region': ['North', 'South', 'North', 'East', 'West']
    }
    
    df = pd.DataFrame(data)
    
    os.makedirs('documents/test', exist_ok=True)
    filepath = 'documents/test/sales_data.xlsx'
    df.to_excel(filepath, index=False)
    
    print(f"✅ Created: {filepath}")
    return filepath

def create_test_html():
    """Create sample HTML file"""
    html_content = """<!DOCTYPE html>
<html>
<head>
    <title>Sample Report</title>
</head>
<body>
    <h1>Q1 Sales Report</h1>
    <p>This is a sample sales report for Q1 2025.</p>
    
    <h2>Key Metrics</h2>
    <ul>
        <li>Total Revenue: $93,100</li>
        <li>Units Sold: 495</li>
        <li>Growth: 15%</li>
    </ul>
    
    <h2>Regional Data</h2>
    <table border="1">
        <tr><th>Region</th><th>Sales</th></tr>
        <tr><td>North</td><td>$58,000</td></tr>
        <tr><td>South</td><td>$3,000</td></tr>
        <tr><td>East</td><td>$22,500</td></tr>
    </table>
    
    <h2>Links</h2>
    <a href="https://example.com/report">Full Report</a>
    <a href="https://example.com/data">Data Source</a>
</body>
</html>"""
    
    os.makedirs('documents/test', exist_ok=True)
    filepath = 'documents/test/report.html'
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"✅ Created: {filepath}")
    return filepath

def create_test_powerpoint():
    """Create sample PowerPoint file"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("❌ python-pptx not installed")
        return
    
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Q1 Business Review"
    title_slide.placeholders[1].text = "January - March 2025"
    
    # Bullet slide
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    bullet_slide.shapes.title.text = "Key Achievements"
    
    content = bullet_slide.shapes.placeholders[1].text_frame
    content.text = "Revenue Growth"
    
    p = content.add_paragraph()
    p.text = "Launched 3 new products"
    p.level = 0
    
    p = content.add_paragraph()
    p.text = "Expanded to 2 new regions"
    p.level = 0
    
    # Summary slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_slide.shapes.title.text = "Next Steps"
    
    content = summary_slide.shapes.placeholders[1].text_frame
    content.text = "Focus on customer retention"
    
    p = content.add_paragraph()
    p.text = "Increase marketing budget by 20%"
    p.level = 0
    
    os.makedirs('documents/test', exist_ok=True)
    filepath = 'documents/test/presentation.pptx'
    prs.save(filepath)
    
    print(f"✅ Created: {filepath}")
    return filepath

def main():
    print("="*50)
    print("Creating Test Office Files")
    print("="*50)
    print()
    
    excel_file = create_test_excel()
    html_file = create_test_html()
    pptx_file = create_test_powerpoint()
    
    print()
    print("="*50)
    print("Test Commands:")
    print("="*50)
    
    if excel_file:
        print("\n📊 Excel:")
        print(f'  python manage_functions.py run excel/process filepath="test/sales_data.xlsx"')
        print(f'  python manage_functions.py run excel/analyze_column filepath="test/sales_data.xlsx" column_name="Revenue"')
    
    if html_file:
        print("\n🌐 HTML:")
        print(f'  python manage_functions.py run html/process filepath="test/report.html" extract="text"')
        print(f'  python manage_functions.py run html/process filepath="test/report.html" extract="links"')
    
    if pptx_file:
        print("\n📊 PowerPoint:")
        print(f'  python manage_functions.py run powerpoint/process filepath="test/presentation.pptx"')
        print(f'  python manage_functions.py run powerpoint/extract_all_text filepath="test/presentation.pptx"')
    
    print()

if __name__ == "__main__":
    main()
