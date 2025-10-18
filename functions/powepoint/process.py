"""
PowerPoint presentation processing and generation functions
Requires: pip install python-pptx
"""

import os
from pathlib import Path

def run(filepath, extract="summary"):
    """
    Process PowerPoint file
    
    Args:
        filepath: Path to PPTX file
        extract: 'summary', 'text', 'notes', 'images'
    
    Returns:
        Extracted content
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"
    
    full_path = os.path.join('documents', filepath) if not os.path.isabs(filepath) else filepath
    
    if not os.path.exists(full_path):
        return f"File not found: {filepath}"
    
    try:
        prs = Presentation(full_path)
        
        if extract == "summary":
            return {
                "total_slides": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
                "slides": [
                    {
                        "slide_num": i+1,
                        "shapes": len(slide.shapes),
                        "title": slide.shapes.title.text if slide.shapes.title else "No title"
                    }
                    for i, slide in enumerate(prs.slides[:10])  # First 10 slides
                ]
            }
        
        elif extract == "text":
            all_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = {
                    "slide_num": i+1,
                    "content": []
                }
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text["content"].append(shape.text)
                all_text.append(slide_text)
            return all_text
        
        elif extract == "notes":
            notes = []
            for i, slide in enumerate(prs.slides):
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        notes.append({
                            "slide_num": i+1,
                            "notes": notes_text
                        })
            return notes if notes else "No speaker notes found"
        
        elif extract == "images":
            images = []
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        images.append({
                            "slide_num": i+1,
                            "image_name": shape.name,
                            "left": shape.left,
                            "top": shape.top
                        })
            return images if images else "No images found"
    
    except Exception as e:
        return f"Error processing PowerPoint: {str(e)}"


def get_slide_text(filepath, slide_num):
    """Get text from specific slide"""
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx not installed"
    
    full_path = os.path.join('documents', filepath) if not os.path.isabs(filepath) else filepath
    
    try:
        prs = Presentation(full_path)
        
        if slide_num < 1 or slide_num > len(prs.slides):
            return f"Invalid slide number. Presentation has {len(prs.slides)} slides"
        
        slide = prs.slides[slide_num - 1]
        text_content = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
        
        return {
            "slide_num": slide_num,
            "text": text_content
        }
    
    except Exception as e:
        return f"Error: {str(e)}"


def create_presentation(title, slides_data, output_path="reports/presentation.pptx"):
    """
    Create PowerPoint presentation
    
    Args:
        title: Presentation title
        slides_data: List of dicts with 'title' and 'content'
        output_path: Where to save
    
    Example:
        slides = [
            {"title": "Introduction", "content": ["Point 1", "Point 2"]},
            {"title": "Overview", "content": ["Detail A", "Detail B"]}
        ]
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "Error: python-pptx not installed"
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        
        # Content slides
        bullet_slide_layout = prs.slide_layouts[1]
        
        for slide_data in slides_data:
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Title
            slide.shapes.title.text = slide_data.get('title', 'Slide')
            
            # Content
            if 'content' in slide_data:
                content_shape = slide.shapes.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                for item in slide_data['content']:
                    p = text_frame.add_paragraph()
                    p.text = str(item)
                    p.level = 0
        
        # Save
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        
        return f"Presentation created: {output_path} ({len(slides_data)+1} slides)"
    
    except Exception as e:
        return f"Error creating presentation: {str(e)}"


def extract_all_text(filepath):
    """Extract all text content from presentation"""
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx not installed"
    
    full_path = os.path.join('documents', filepath) if not os.path.isabs(filepath) else filepath
    
    try:
        prs = Presentation(full_path)
        all_text = []
        
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    all_text.append(shape.text)
        
        return "\n\n".join(all_text)
    
    except Exception as e:
        return f"Error: {str(e)}"


def count_slides(filepath):
    """Count total slides in presentation"""
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx not installed"
    
    full_path = os.path.join('documents', filepath) if not os.path.isabs(filepath) else filepath
    
    try:
        prs = Presentation(full_path)
        return len(prs.slides)
    except Exception as e:
        return f"Error: {str(e)}"


def search_in_presentation(filepath, search_term):
    """Search for term in presentation"""
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx not installed"
    
    full_path = os.path.join('documents', filepath) if not os.path.isabs(filepath) else filepath
    
    try:
        prs = Presentation(full_path)
        matches = []
        
        for i, slide in enumerate(prs.slides):
            slide_matches = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and search_term.lower() in shape.text.lower():
                    slide_matches.append(shape.text)
            
            if slide_matches:
                matches.append({
                    "slide_num": i+1,
                    "matches": slide_matches
                })
        
        return matches if matches else f"No matches found for '{search_term}'"
    
    except Exception as e:
        return f"Error: {str(e)}"


def pptx_to_pdf(filepath, output_path=None):
    """
    Convert PowerPoint to PDF (requires LibreOffice or MS Office)
    
    Note: This is a placeholder. Actual conversion requires:
    - Linux/Mac: LibreOffice (libreoffice --headless --convert-to pdf file.pptx)
    - Windows: comtypes library + MS Office
    """
    import subprocess
    import platform
    
    full_path = os.path.join('documents', filepath) if not os.path.isabs(filepath) else filepath
    
    if output_path is None:
        output_path = full_path.replace('.pptx', '.pdf')
    
    try:
        if platform.system() in ['Linux', 'Darwin']:  # Linux or Mac
            result = subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', 
                 os.path.dirname(output_path), full_path],
                capture_output=True,
                timeout=60
            )
            if result.returncode == 0:
                return f"PDF created: {output_path}"
            else:
                return "Error: LibreOffice not found or conversion failed"
        else:
            return "PDF conversion requires LibreOffice (Linux/Mac) or MS Office (Windows with comtypes)"
    
    except Exception as e:
        return f"Error: {str(e)}"
