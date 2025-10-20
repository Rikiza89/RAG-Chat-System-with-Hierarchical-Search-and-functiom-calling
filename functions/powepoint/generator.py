"""
PowerPoint File Generator - Create custom presentations from user requests
Supports various templates and customization options
"""

import os
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def run(content, title="Generated Presentation", template="basic", output_path=None, **kwargs):
    """
    Generate PowerPoint presentation based on user request
    
    Args:
        content: Main content (text, list, dict, or structured data)
        title: Presentation title
        template: Template type ('basic', 'business', 'pitch', 'report', 'academic', 'minimal')
        output_path: Where to save (default: reports/generated_TIMESTAMP.pptx)
        **kwargs: Additional customization options
            - theme: 'light' or 'dark' (default: 'light')
            - color: Primary color as hex (default: '#667eea')
            - slides: List of slide dictionaries for multi-slide presentations
            - author: Author name (default: 'RAG System')
            - subtitle: Subtitle for title slide
    
    Returns:
        Path to generated PowerPoint file
    
    Examples:
        # Simple presentation
        <run:pptx/generator content="Welcome to our presentation" 
                            title="Company Overview">
        
        # Business report
        <run:pptx/generator content='{"Q1 Revenue": "$1.2M", "Growth": "15%"}' 
                           title="Q1 Results" 
                           template="business">
        
        # Multi-slide presentation
        <run:pptx/generator slides='[{"title": "Intro", "content": "..."}, 
                                     {"title": "Data", "content": "..."}]'
                           title="Annual Report"
                           template="report">
    """
    
    # Set default output path
    if output_path is None:
        os.makedirs('reports', exist_ok=True)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_path = f'reports/{title.lower().replace(" ", "_")}_{timestamp}.pptx'
    else:
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    
    # Get template generator
    templates = {
        'basic': generate_basic_presentation,
        'business': generate_business_presentation,
        'pitch': generate_pitch_presentation,
        'report': generate_report_presentation,
        'academic': generate_academic_presentation,
        'minimal': generate_minimal_presentation
    }
    
    generator = templates.get(template, generate_basic_presentation)
    
    # Generate presentation
    prs = generator(content, title, **kwargs)
    
    # Save file
    prs.save(output_path)
    
    return f"✅ PowerPoint file created: {output_path}"


def generate_basic_presentation(content, title, theme='light', color='#667eea', 
                                subtitle='', author='RAG System', slides=None, **kwargs):
    """Generate basic presentation template"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Parse color
    rgb_color = hex_to_rgb(color)
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_title_slide(title_slide, title, subtitle, rgb_color, theme)
    
    # Content slides
    if slides:
        if isinstance(slides, str):
            try:
                slides = json.loads(slides)
            except:
                slides = [{"title": "Content", "content": slides}]
        
        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_content_slide(slide, slide_data.get('title', 'Slide'), 
                            slide_data.get('content', ''), rgb_color, theme)
    else:
        # Single content slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_content_slide(slide, "Content", content, rgb_color, theme)
    
    # Thank you slide
    thank_you_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_thank_you_slide(thank_you_slide, rgb_color, theme, author)
    
    return prs


def generate_business_presentation(content, title, color='#2E86AB', **kwargs):
    """Generate professional business presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    rgb_color = hex_to_rgb(color)
    
    # Title slide with company branding
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_business_title_slide(title_slide, title, kwargs.get('subtitle', ''), rgb_color)
    
    # Executive summary
    if isinstance(content, dict):
        summary_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_executive_summary_slide(summary_slide, content, rgb_color)
        
        # Detail slides for each key
        for key, value in content.items():
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_detail_slide(slide, key, value, rgb_color)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_content_slide(slide, "Overview", content, rgb_color, 'light')
    
    # Closing slide
    closing_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_business_closing_slide(closing_slide, rgb_color)
    
    return prs


def generate_pitch_presentation(content, title, color='#FF6B35', **kwargs):
    """Generate startup pitch deck presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    rgb_color = hex_to_rgb(color)
    
    # Cover slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_pitch_title_slide(title_slide, title, kwargs.get('subtitle', ''), rgb_color)
    
    # Problem slide
    problem_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_slide(problem_slide, "The Problem", 
                     "Every great solution starts with understanding the problem", rgb_color)
    
    # Solution slide
    solution_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(solution_slide, "Our Solution", content, rgb_color, 'light')
    
    # Market opportunity
    market_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_slide(market_slide, "Market Opportunity", 
                     "A massive market ready for disruption", rgb_color)
    
    # Contact slide
    contact_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_contact_slide(contact_slide, rgb_color)
    
    return prs


def generate_report_presentation(content, title, color='#4A5859', **kwargs):
    """Generate data-focused report presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    rgb_color = hex_to_rgb(color)
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_report_title_slide(title_slide, title, rgb_color)
    
    # Key findings
    if isinstance(content, dict):
        findings_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_key_findings_slide(findings_slide, content, rgb_color)
        
        # Individual metric slides
        for key, value in list(content.items())[:5]:  # Limit to 5 metrics
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_metric_slide(slide, key, value, rgb_color)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_content_slide(slide, "Report Summary", content, rgb_color, 'light')
    
    # Summary slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_summary_slide(summary_slide, rgb_color)
    
    return prs


def generate_academic_presentation(content, title, color='#1B4965', **kwargs):
    """Generate academic/research presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    rgb_color = hex_to_rgb(color)
    
    # Title slide with author
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_academic_title_slide(title_slide, title, kwargs.get('author', 'Researcher'), rgb_color)
    
    # Abstract
    abstract_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(abstract_slide, "Abstract", content, rgb_color, 'light')
    
    # Methodology (if dict provided)
    if isinstance(content, dict):
        method_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_content_slide(method_slide, "Methodology", 
                         format_dict_content(content), rgb_color, 'light')
    
    # References slide
    ref_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_references_slide(ref_slide, rgb_color)
    
    return prs


def generate_minimal_presentation(content, title, color='#000000', **kwargs):
    """Generate minimal/clean presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    rgb_color = hex_to_rgb(color)
    
    # Minimal title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_minimal_title_slide(title_slide, title, rgb_color)
    
    # Content slides
    slides = kwargs.get('slides')
    if slides:
        if isinstance(slides, str):
            try:
                slides = json.loads(slides)
            except:
                slides = [{"title": title, "content": content}]
        
        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_minimal_content_slide(slide, slide_data.get('title', ''), 
                                     slide_data.get('content', ''), rgb_color)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_minimal_content_slide(slide, "", content, rgb_color)
    
    return prs


# Slide builder functions

def add_title_slide(slide, title, subtitle, rgb_color, theme):
    """Add formatted title slide"""
    bg_color = RGBColor(255, 255, 255) if theme == 'light' else RGBColor(26, 26, 26)
    text_color = RGBColor(51, 51, 51) if theme == 'light' else RGBColor(224, 224, 224)
    
    # Background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = bg_color
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_color
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = text_color
        subtitle_para.alignment = PP_ALIGN.CENTER


def add_content_slide(slide, slide_title, content, rgb_color, theme):
    """Add formatted content slide"""
    bg_color = RGBColor(255, 255, 255) if theme == 'light' else RGBColor(26, 26, 26)
    text_color = RGBColor(51, 51, 51) if theme == 'light' else RGBColor(224, 224, 224)
    
    # Background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = bg_color
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = slide_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_color
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    formatted_content = format_content_for_slide(content)
    content_frame.text = formatted_content
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = text_color
        paragraph.space_after = Pt(12)


def add_thank_you_slide(slide, rgb_color, theme, author):
    """Add thank you slide"""
    bg_color = RGBColor(255, 255, 255) if theme == 'light' else RGBColor(26, 26, 26)
    
    # Background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = bg_color
    background.line.fill.background()
    
    # Thank you text
    thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(60)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = rgb_color
    thank_you_para.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = f"Generated by {author} | {datetime.now().strftime('%B %d, %Y')}"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = RGBColor(153, 153, 153)
    footer_para.alignment = PP_ALIGN.CENTER


def add_business_title_slide(slide, title, subtitle, rgb_color):
    """Add business-style title slide"""
    # Gradient background effect (solid color as approximation)
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = rgb_color
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER


def add_executive_summary_slide(slide, data, rgb_color):
    """Add executive summary slide with key metrics"""
    # Background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Executive Summary"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_color
    
    # Metrics grid
    y_pos = 2
    for i, (key, value) in enumerate(list(data.items())[:4]):
        x_pos = 1 if i % 2 == 0 else 5.5
        if i >= 2:
            y_pos = 4.5
        
        metric_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos if i < 2 else 4.5), 
                                             Inches(3.5), Inches(1.5))
        metric_frame = metric_box.text_frame
        
        # Key
        metric_frame.text = key.replace('_', ' ').title()
        key_para = metric_frame.paragraphs[0]
        key_para.font.size = Pt(20)
        key_para.font.color.rgb = RGBColor(102, 102, 102)
        
        # Value
        value_para = metric_frame.add_paragraph()
        value_para.text = str(value)
        value_para.font.size = Pt(36)
        value_para.font.bold = True
        value_para.font.color.rgb = rgb_color


def add_detail_slide(slide, key, value, rgb_color):
    """Add detail slide for specific metric"""
    add_content_slide(slide, key.replace('_', ' ').title(), str(value), rgb_color, 'light')


def add_business_closing_slide(slide, rgb_color):
    """Add business closing slide"""
    # Background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = rgb_color
    background.line.fill.background()
    
    # Text
    text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    text_frame = text_box.text_frame
    text_frame.text = "Questions?"
    text_para = text_frame.paragraphs[0]
    text_para.font.size = Pt(60)
    text_para.font.bold = True
    text_para.font.color.rgb = RGBColor(255, 255, 255)
    text_para.alignment = PP_ALIGN.CENTER


def add_pitch_title_slide(slide, title, subtitle, rgb_color):
    """Add pitch deck title slide"""
    add_business_title_slide(slide, title, subtitle, rgb_color)


def add_section_slide(slide, section_title, description, rgb_color):
    """Add section divider slide"""
    # Background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 245, 245)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_color
    title_para.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(24)
    desc_para.font.color.rgb = RGBColor(102, 102, 102)
    desc_para.alignment = PP_ALIGN.CENTER


def add_contact_slide(slide, rgb_color):
    """Add contact slide"""
    add_section_slide(slide, "Let's Connect", "Reach out to discuss opportunities", rgb_color)


def add_report_title_slide(slide, title, rgb_color):
    """Add report title slide"""
    add_business_title_slide(slide, title, f"Report Generated: {datetime.now().strftime('%B %d, %Y')}", 
                           rgb_color)


def add_key_findings_slide(slide, data, rgb_color):
    """Add key findings slide"""
    add_executive_summary_slide(slide, data, rgb_color)


def add_metric_slide(slide, metric_name, metric_value, rgb_color):
    """Add individual metric slide"""
    # Background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Metric name
    name_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    name_frame = name_box.text_frame
    name_frame.text = metric_name.replace('_', ' ').title()
    name_para = name_frame.paragraphs[0]
    name_para.font.size = Pt(36)
    name_para.font.color.rgb = RGBColor(102, 102, 102)
    name_para.alignment = PP_ALIGN.CENTER
    
    # Metric value
    value_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.5))
    value_frame = value_box.text_frame
    value_frame.text = str(metric_value)
    value_para = value_frame.paragraphs[0]
    value_para.font.size = Pt(72)
    value_para.font.bold = True
    value_para.font.color.rgb = rgb_color
    value_para.alignment = PP_ALIGN.CENTER


def add_summary_slide(slide, rgb_color):
    """Add summary slide"""
    add_section_slide(slide, "Summary", "Key takeaways from this report", rgb_color)


def add_academic_title_slide(slide, title, author, rgb_color):
    """Add academic title slide"""
    # Background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_color
    title_para.alignment = PP_ALIGN.CENTER
    
    # Author
    author_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
    author_frame = author_box.text_frame
    author_frame.text = author
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(24)
    author_para.font.color.rgb = RGBColor(102, 102, 102)
    author_para.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime('%B %d, %Y')
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(153, 153, 153)
    date_para.alignment = PP_ALIGN.CENTER


def add_references_slide(slide, rgb_color):
    """Add references slide"""
    add_content_slide(slide, "References", "1. Reference materials here\n2. Additional sources", 
                     rgb_color, 'light')


def add_minimal_title_slide(slide, title, rgb_color):
    """Add minimal title slide"""
    # White background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_color
    title_para.alignment = PP_ALIGN.LEFT


def add_minimal_content_slide(slide, slide_title, content, rgb_color):
    """Add minimal content slide"""
    # White background
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Title (optional)
    y_start = 1.5
    if slide_title:
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = rgb_color
        y_start = 2.2
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(y_start), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    formatted_content = format_content_for_slide(content)
    content_frame.text = formatted_content
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)


# Helper functions

def hex_to_rgb(hex_color):
    """Convert hex color to RGB"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def format_content_for_slide(content):
    """Format content for slide display"""
    if isinstance(content, dict):
        return format_dict_content(content)
    elif isinstance(content, list):
        return format_list_content(content)
    else:
        return str(content)


def format_dict_content(data):
    """Format dictionary as bullet points"""
    lines = []
    for key, value in data.items():
        lines.append(f"• {key.replace('_', ' ').title()}: {value}")
    return '\n'.join(lines)


def format_list_content(data):
    """Format list as bullet points"""
    lines = []
    for item in data:
        lines.append(f"• {item}")
    return '\n'.join(lines)
